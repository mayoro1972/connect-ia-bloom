"""
Deck TransferAI — Template de référence
Version épurée, aérée, adaptable par prospect et secteur d'activité.
Usage : modifier la section DONNÉES PROSPECT puis exécuter.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import subprocess

# ── Palette exacte ────────────────────────────────────────────────────────
BLEU_FONCE  = RGBColor(0x10, 0x26, 0x3F)
BLEU_MOYEN  = RGBColor(0x16, 0x35, 0x56)
BLEU_CLAIR  = RGBColor(0x24, 0x46, 0x67)
BLEU_TEXTE  = RGBColor(0x61, 0x70, 0x86)
ORANGE      = RGBColor(0xE7, 0x6F, 0x1D)
OR_AMBRE    = RGBColor(0xC8, 0x8C, 0x3A)
TEAL        = RGBColor(0x1C, 0x8A, 0x78)
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
CREME       = RGBColor(0xF5, 0xEF, 0xE7)
BLEU_PALE   = RGBColor(0xCF, 0xE0, 0xF0)
BLEU_GRIS   = RGBColor(0xC9, 0xD4, 0xE1)


# ════════════════════════════════════════════════════════════════════════
# ▶  DONNÉES PROSPECT — modifier ici pour chaque nouveau deck
# ════════════════════════════════════════════════════════════════════════

PROSPECT        = "Orange Côte d'Ivoire"
SECTEUR         = "Télécommunications"
PAYS            = "Côte d'Ivoire"

# Phrase d'accroche — 1 ligne, spécifique au secteur
ACCROCHE        = "Transformer vos flux clients et vos processus internes avec l'IA"

# 3 enjeux identifiés lors du pré-audit (phrases courtes)
ENJEUX = [
    "Améliorer la qualité et la cohérence des réponses service client",
    "Accélérer la production des reportings managériaux",
    "Structurer et diffuser les procédures et scripts terrain",
]

# 3 cas d'usage prioritaires (icône, titre court, gain attendu)
CAS_USAGE = [
    ("💬", "Copilote service client",      "Réponses homogènes, délai réduit, continuité multicanal"),
    ("📋", "Synthèse tickets et pilotage", "Reporting automatisé, alertes lisibles, meilleur arbitrage"),
    ("📖", "Base de connaissances terrain","Scripts unifiés, onboarding accéléré, moins d'erreurs"),
]

# 3 gains mesurables du pilote
GAINS = [
    ("Délai de réponse",       "− 30 à 50 %",  TEAL),
    ("Temps de reporting",     "− 60 %",        ORANGE),
    ("Qualité des réponses",   "+ homogénéité", BLEU_CLAIR),
]

# Plan 90 jours — 4 étapes (numéro, délai, action)
PLAN = [
    ("01", "J+0",  "Audit & cadrage"),
    ("02", "J+15", "Premier pilote"),
    ("03", "J+45", "Formation & déploiement"),
    ("04", "J+90", "Mesure & extension"),
]

# URLs
PACK_ID      = "pack-1780665567601-jp48u6z4"
AUDIT_URL    = f"https://www.transferai.ci/questionnaire-audit?pack_id={PACK_ID}"
CALENDLY_URL = "https://calendly.com/contact-transferai/30min"

# Équipe — fixe pour tous les decks
EQUIPE = [
    ("Casimir Kassi Beda",  "Directeur Général",                              "contact@transferai.ci"),
    ("Marius Ayoro",        "Directeur Développement & Partenariats",         "marius.ayoro@transferai.ci"),
    ("Soulemane Konate",    "Directeur IA & Innovation",                      "soulemane.konate@transferai.ci"),
    ("Eric N'Guessan",      "Directeur Audit, Pédagogie & Certification",     "eric.nguessan@transferai.ci"),
    ("Médard Séry",         "Consultant expert · Data & plateformes IA",      "medard.sery@transferai.ci"),
]

OUT_PPTX = f"/Users/marius_ayoro/Downloads/Deck_TransferAI_{PROSPECT.replace(' ', '_').replace(chr(39), '')}_V3.pptx"


# ════════════════════════════════════════════════════════════════════════
# HELPERS
# ════════════════════════════════════════════════════════════════════════
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s

def tx(slide, l, t, w, h, text, size=11, bold=False, color=BLANC,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def footer(slide, num):
    rect(slide, 0.5, 7.02, 12.33, 0.02, fill=ORANGE)
    tx(slide, 0.6, 7.06, 10.0, 0.22,
       "TransferAI  ·  Hub IA de NettelecomCI  ·  contact@transferai.ci",
       size=7, color=BLEU_TEXTE)
    tx(slide, 12.2, 7.04, 0.80, 0.24,
       f"{num:02d}", size=9, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════

def slide_couverture(prs):
    sl = blank(prs)
    # Fond
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 8.60, 0, 4.733, 7.5, fill=BLEU_MOYEN)
    rect(sl, 0.55, 0.85, 0.05, 5.50, fill=ORANGE)

    # Contenu gauche
    tx(sl, 0.88, 0.82, 4.0, 0.26, "PRÉSENTATION COMMERCIALE",
       size=7.5, bold=True, color=OR_AMBRE)
    tx(sl, 0.90, 1.28, 7.40, 0.72,
       f"TransferAI × {PROSPECT}",
       size=24, bold=True, color=BLANC)
    tx(sl, 0.90, 2.16, 7.30, 0.38, ACCROCHE,
       size=13, color=BLEU_PALE, italic=True)

    # Séparateur
    rect(sl, 0.90, 2.76, 3.20, 0.03, fill=ORANGE)

    # Secteur + pays
    tx(sl, 0.90, 2.98, 7.0, 0.28,
       f"{SECTEUR}  ·  {PAYS}",
       size=10, color=BLEU_GRIS)

    # Enjeux résumés
    for i, e in enumerate(ENJEUX):
        tx(sl, 0.90, 3.55 + i * 0.52, 7.20, 0.42,
           f"▸  {e}", size=10.5, color=BLANC)

    # Colonne droite
    tx(sl, 8.90, 1.20, 3.80, 0.28, "TransferAI Africa",
       size=9, bold=True, color=OR_AMBRE)
    tx(sl, 8.90, 1.62, 3.70, 0.26, "Hub IA de NettelecomCI",
       size=10.5, bold=True, color=BLANC)
    tx(sl, 8.90, 2.06, 3.70, 1.00,
       "Présence locale en Côte d'Ivoire · Expertise internationale · "
       "Audit, formation et accompagnement 90 jours.",
       size=9.5, color=BLEU_PALE, wrap=True)

    rect(sl, 8.90, 3.22, 3.70, 0.02, fill=BLEU_GRIS)

    tx(sl, 8.90, 3.40, 3.70, 0.26, "Rendez-vous →",
       size=9, bold=True, color=ORANGE)
    tx(sl, 8.90, 3.76, 3.70, 0.28, CALENDLY_URL,
       size=8, color=BLEU_PALE, italic=True)

    tx(sl, 8.90, 4.24, 3.70, 0.26, "Formulaire audit →",
       size=9, bold=True, color=ORANGE)
    tx(sl, 8.90, 4.60, 3.70, 0.28, "transferai.ci/questionnaire-audit",
       size=8, color=BLEU_PALE, italic=True)

    footer(sl, 1)


def slide_enjeux(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)

    # En-tête
    rect(sl, 0, 0, 13.333, 1.40, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.30, 10.0, 0.50,
       "Enjeux identifiés lors du pré-audit",
       size=20, bold=True, color=BLANC)
    tx(sl, 0.74, 0.88, 9.0, 0.30,
       f"{PROSPECT}  ·  {SECTEUR}",
       size=9.5, color=BLEU_PALE)

    # 3 cartes enjeux
    for i, enjeu in enumerate(ENJEUX):
        lx = 0.72 + i * 4.22
        rect(sl, lx, 1.75, 3.90, 3.20, fill=CREME)
        rect(sl, lx, 1.75, 3.90, 0.07, fill=ORANGE)
        # Numéro
        tx(sl, lx + 0.20, 1.90, 0.60, 0.55,
           f"0{i+1}", size=26, bold=True, color=ORANGE)
        # Texte enjeu
        tx(sl, lx + 0.20, 2.62, 3.50, 1.80,
           enjeu, size=13, bold=True, color=BLEU_FONCE, wrap=True)

    # Message bas de page
    rect(sl, 0.72, 5.20, 11.89, 0.82, fill=BLEU_FONCE)
    tx(sl, 1.00, 5.42, 11.40, 0.42,
       "Ces enjeux sont les points d'entrée prioritaires pour un premier déploiement IA ciblé, "
       "mesurable et défendable dès les 90 premiers jours.",
       size=10.5, color=BLANC, italic=True, wrap=True)

    footer(sl, 2)


def slide_cas_usage(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)

    rect(sl, 0, 0, 13.333, 1.40, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.30, 10.0, 0.50,
       "Cas d'usage prioritaires",
       size=20, bold=True, color=BLANC)
    tx(sl, 0.74, 0.88, 9.0, 0.30,
       "Trois leviers de transformation IA adaptés à votre secteur",
       size=9.5, color=BLEU_PALE)

    for i, (ico, titre, gain) in enumerate(CAS_USAGE):
        lx = 0.72 + i * 4.22
        # Carte fond
        rect(sl, lx, 1.72, 3.90, 4.60, fill=CREME)
        rect(sl, lx, 1.72, 3.90, 0.07, fill=BLEU_FONCE)

        # Icône
        rect(sl, lx + 0.20, 1.95, 0.70, 0.70, fill=BLEU_FONCE)
        tx(sl, lx + 0.20, 1.92, 0.70, 0.70,
           ico, size=22, color=BLANC, align=PP_ALIGN.CENTER)

        # Titre
        tx(sl, lx + 1.05, 2.05, 2.65, 0.55,
           titre, size=13, bold=True, color=BLEU_FONCE, wrap=True)

        # Séparateur
        rect(sl, lx + 0.20, 2.82, 3.50, 0.02, fill=BLEU_GRIS)

        # Gain
        rect(sl, lx + 0.20, 3.06, 0.65, 0.24, fill=TEAL)
        tx(sl, lx + 0.20, 3.06, 0.65, 0.24,
           "GAIN", size=7, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        tx(sl, lx + 0.20, 3.42, 3.50, 1.60,
           gain, size=11, color=BLEU_TEXTE, wrap=True)

    footer(sl, 3)


def slide_roi(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)

    # Titre
    tx(sl, 0.72, 0.65, 10.0, 0.55,
       "Résultats attendus dès le pilote",
       size=22, bold=True, color=BLANC)
    tx(sl, 0.74, 1.30, 9.0, 0.30,
       "Des gains visibles, mesurés et défendables",
       size=10.5, color=BLEU_PALE, italic=True)

    # 3 KPI cards
    for i, (label, valeur, couleur) in enumerate(GAINS):
        lx = 0.72 + i * 4.22
        rect(sl, lx, 2.00, 3.90, 2.20, fill=BLEU_MOYEN)
        rect(sl, lx, 2.00, 3.90, 0.07, fill=couleur)
        tx(sl, lx + 0.24, 2.28, 3.42, 0.40,
           label, size=10, color=BLEU_PALE)
        tx(sl, lx + 0.24, 2.80, 3.42, 0.70,
           valeur, size=28, bold=True, color=couleur, wrap=False)

    # Gouvernance
    rect(sl, 0.72, 4.52, 11.89, 1.60, fill=BLEU_MOYEN)
    tx(sl, 1.00, 4.75, 3.0, 0.30,
       "Gouvernance & supervision", size=9, bold=True, color=OR_AMBRE)
    principes = [
        "Supervision humaine de tous les usages IA",
        "KPI validés avant toute extension",
        "Pilotage simple, sobre et explicable",
    ]
    for i, p in enumerate(principes):
        lx = 1.00 + i * 3.96
        tx(sl, lx, 5.15, 3.70, 0.70,
           f"▸  {p}", size=10, color=BLANC, wrap=True)

    footer(sl, 4)


def slide_plan(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)

    rect(sl, 0, 0, 13.333, 1.40, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.30, 10.0, 0.50,
       "Plan d'action 90 jours",
       size=20, bold=True, color=BLANC)
    tx(sl, 0.74, 0.88, 9.0, 0.30,
       "Audit · Formation · Déploiement · Extension",
       size=9.5, color=BLEU_PALE)

    # 4 étapes
    for i, (num, delai, action) in enumerate(PLAN):
        lx = 0.72 + i * 3.15
        rect(sl, lx, 1.72, 2.90, 2.30, fill=CREME)
        rect(sl, lx, 1.72, 2.90, 0.07, fill=BLEU_FONCE)
        # Numéro
        tx(sl, lx + 0.20, 1.88, 0.70, 0.50,
           num, size=26, bold=True, color=BLEU_FONCE, wrap=False)
        # Délai
        tx(sl, lx + 1.00, 1.96, 1.70, 0.34,
           delai, size=13, bold=True, color=ORANGE, wrap=False)
        # Séparateur
        rect(sl, lx + 0.20, 2.52, 2.50, 0.02, fill=BLEU_GRIS)
        # Action
        tx(sl, lx + 0.20, 2.66, 2.50, 0.90,
           action, size=12, bold=True, color=BLEU_FONCE, wrap=True)

    # Accompagnement bas
    rect(sl, 0.72, 4.26, 5.75, 2.00, fill=BLEU_FONCE)
    tx(sl, 0.96, 4.48, 5.20, 0.30, "Accompagnement TransferAI",
       size=9, bold=True, color=OR_AMBRE)
    accomp = [
        "Points de suivi hebdomadaires sur l'usage réel",
        "Ajustement des prompts et procédures",
        "Revue des KPI et décision d'extension",
    ]
    for i, a in enumerate(accomp):
        tx(sl, 0.96, 4.90 + i * 0.38, 5.20, 0.34,
           f"·  {a}", size=10, color=BLANC)

    # Formation bas
    rect(sl, 6.72, 4.26, 5.89, 2.00, fill=CREME)
    tx(sl, 6.96, 4.48, 5.40, 0.30, "Formation ciblée",
       size=9, bold=True, color=BLEU_FONCE)
    form = [
        "Usage du copilote en situation réelle",
        "Validation et gouvernance des réponses IA",
        "Mise à jour des scripts et procédures",
    ]
    for i, f in enumerate(form):
        tx(sl, 6.96, 4.90 + i * 0.38, 5.40, 0.34,
           f"·  {f}", size=10, color=BLEU_TEXTE)

    footer(sl, 5)


def slide_cta(prs):
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 7.80, 0, 5.533, 7.5, fill=BLEU_MOYEN)
    rect(sl, 0.55, 0.85, 0.05, 5.50, fill=ORANGE)

    # Gauche — CTA
    tx(sl, 0.88, 0.90, 4.0, 0.28, "PROCHAINE ÉTAPE",
       size=7.5, bold=True, color=OR_AMBRE)
    tx(sl, 0.88, 1.35, 6.60, 0.85,
       "Planifier un audit stratégique gratuit\net un rendez-vous expert de 30 minutes",
       size=19, bold=True, color=BLANC, wrap=True)

    rect(sl, 0.88, 2.50, 6.55, 0.02, fill=ORANGE)

    # Liens
    rect(sl, 0.88, 2.72, 6.55, 0.64, fill=BLEU_CLAIR)
    tx(sl, 1.08, 2.84, 1.50, 0.26, "Rendez-vous",
       size=8.5, bold=True, color=OR_AMBRE)
    tx(sl, 2.72, 2.88, 4.50, 0.26, CALENDLY_URL,
       size=8.5, color=BLEU_PALE, italic=True)

    rect(sl, 0.88, 3.50, 6.55, 0.64, fill=BLEU_CLAIR)
    tx(sl, 1.08, 3.62, 1.70, 0.26, "Formulaire audit",
       size=8.5, bold=True, color=OR_AMBRE)
    tx(sl, 2.72, 3.66, 4.50, 0.26, "transferai.ci/questionnaire-audit",
       size=8.5, color=BLEU_PALE, italic=True)

    tx(sl, 0.88, 4.44, 6.55, 0.80,
       "Un premier échange de 30 minutes pour cadrer les priorités "
       "et définir ensemble le périmètre du pilote.",
       size=10.5, color=BLANC, italic=True, wrap=True)

    # Droite — Équipe
    tx(sl, 8.10, 1.00, 4.80, 0.30, "Équipe dirigeante",
       size=9, bold=True, color=OR_AMBRE)

    for i, (nom, role, email) in enumerate(EQUIPE):
        y = 1.50 + i * 1.06
        rect(sl, 8.10, y, 4.80, 0.02, fill=BLEU_CLAIR)
        tx(sl, 8.10, y + 0.10, 4.80, 0.28,
           nom, size=10.5, bold=True, color=BLANC)
        tx(sl, 8.10, y + 0.40, 4.80, 0.26,
           role, size=8.5, color=BLEU_PALE, wrap=True)
        tx(sl, 8.10, y + 0.68, 4.80, 0.22,
           email, size=7.5, color=ORANGE, italic=True)

    footer(sl, 6)


# ════════════════════════════════════════════════════════════════════════
# BUILD + EXPORT
# ════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_couverture(prs)
    slide_enjeux(prs)
    slide_cas_usage(prs)
    slide_roi(prs)
    slide_plan(prs)
    slide_cta(prs)
    prs.save(OUT_PPTX)
    print(f"✓ PPTX : {OUT_PPTX}")
    return OUT_PPTX

def export_pdf(pptx_path):
    pdf = pptx_path.replace(".pptx", ".pdf")
    script = f'''
tell application "Keynote"
    set d to open POSIX file "{pptx_path}"
    delay 4
    export d to POSIX file "{pdf}" as PDF with properties {{PDF image quality:Best}}
    close d saving no
end tell
'''
    r = subprocess.run(["osascript", "-e", script],
                       capture_output=True, text=True, timeout=90)
    if r.returncode == 0:
        print(f"✓ PDF  : {pdf}")
    else:
        print(f"⚠ PDF : {r.stderr.strip()}")
    return pdf

if __name__ == "__main__":
    pptx = build()
    export_pdf(pptx)
    print("✓ Terminé — fichiers dans ~/Downloads/")
