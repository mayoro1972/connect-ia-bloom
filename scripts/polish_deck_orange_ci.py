"""
Deck TransferAI x Orange CI — Version V2 polie
Couleurs exactes du document original + zero chevauchement
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import subprocess, os

# ── Palette exacte extraite du document original ───────────────────────────
BLEU_FONCE  = RGBColor(0x10, 0x26, 0x3F)  # fond dark principal
BLEU_MOYEN  = RGBColor(0x16, 0x35, 0x56)  # fond section droite
BLEU_CLAIR  = RGBColor(0x24, 0x46, 0x67)  # cards secondaires
BLEU_TEXTE  = RGBColor(0x61, 0x70, 0x86)  # texte gris-bleu
ORANGE      = RGBColor(0xE7, 0x6F, 0x1D)  # orange CI accent
OR_AMBRE    = RGBColor(0xC8, 0x8C, 0x3A)  # or / section label
TEAL        = RGBColor(0x1C, 0x8A, 0x78)  # vert KPI
BLANC       = RGBColor(0xFF, 0xFF, 0xFF)
CREME       = RGBColor(0xF5, 0xEF, 0xE7)  # fond clair
BLEU_PALE   = RGBColor(0xCF, 0xE0, 0xF0)  # texte secondaire clair
BLEU_GRIS   = RGBColor(0xC9, 0xD4, 0xE1)  # séparateurs

OUT_PPTX = "/Users/marius_ayoro/Downloads/Deck_TransferAI_Orange_CI_V2.pptx"
W = Inches(13.333)
H = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════════════════
def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s


def txb(slide, l, t, w, h, text, size=10, bold=False, color=BLANC,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    """Ajoute un text box avec une seule ligne de texte, sans chevauchement."""
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tf


def txb_lines(slide, l, t, w, h, lines):
    """
    lines = liste de (text, size, bold, color, align)
    Chaque ligne est un paragraphe séparé — aucun chevauchement possible.
    """
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.text_frame.word_wrap = True
    first = True
    for (text, size, bold, color, align) in lines:
        p = tf.text_frame.paragraphs[0] if first else tf.text_frame.add_paragraph()
        first = False
        p.alignment = align or PP_ALIGN.LEFT
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bool(bold)
        run.font.color.rgb = color or BLANC
    return tf


def footer(slide, num):
    rect(slide, 0.50, 6.84, 12.33, 0.03, fill=ORANGE)
    txb(slide, 0.60, 6.90, 10.0, 0.22,
        "TransferAI  |  Hub IA de NettelecomCI  |  Modèle commercial premium",
        size=7.5, color=BLEU_TEXTE)
    txb(slide, 12.40, 6.88, 0.65, 0.24,
        f"{num:02d}", size=9, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)


def pill_shape(slide, l, t, w, h, text, bg, fg=BLANC, size=9):
    r = rect(slide, l, t, w, h, fill=bg)
    tf = r.text_frame
    tf.word_wrap = False
    tf.margin_top = Inches(h/2 - 0.11)
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = fg
    return r


# ══════════════════════════════════════════════════════════════════════════
# COMPOSANTS SPÉCIALISÉS
# ══════════════════════════════════════════════════════════════════════════
def agenda_row(slide, num, titre, sous, y):
    rect(slide, 0.95, y, 6.30, 0.58, fill=CREME)
    rect(slide, 0.95, y, 0.55, 0.58, fill=BLEU_FONCE)
    txb(slide, 0.95, y+0.12, 0.55, 0.34, num,
        size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    txb(slide, 1.62, y+0.07, 2.50, 0.26, titre,
        size=11, bold=True, color=BLEU_FONCE)
    txb(slide, 4.20, y+0.09, 3.00, 0.24, sous,
        size=9, color=BLEU_TEXTE)


def use_case_card(slide, l, t, w, icon, titre, avant, apres, kpi):
    """Carte cas d'usage — hauteur fixe 3.90", aucun texte superposé."""
    h = 3.90
    rect(slide, l, t, w, h, fill=CREME)
    # Barre top
    rect(slide, l, t, w, 0.06, fill=ORANGE)
    # Icone
    rect(slide, l+0.18, t+0.18, 0.44, 0.44, fill=BLEU_FONCE)
    txb(slide, l+0.18, t+0.16, 0.44, 0.44,
        icon, size=18, color=BLANC, align=PP_ALIGN.CENTER)
    # Titre
    txb(slide, l+0.72, t+0.22, w-0.90, 0.40,
        titre, size=11, bold=True, color=BLEU_FONCE, wrap=True)
    # Avant
    rect(slide, l+0.18, t+0.76, 0.55, 0.22, fill=ORANGE)
    txb(slide, l+0.18, t+0.76, 0.55, 0.22,
        "AVANT", size=7.5, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    txb(slide, l+0.18, t+1.02, w-0.36, 0.65,
        avant, size=9.5, color=BLEU_TEXTE, wrap=True)
    # Après
    rect(slide, l+0.18, t+1.74, 0.55, 0.22, fill=TEAL)
    txb(slide, l+0.18, t+1.74, 0.55, 0.22,
        "APRÈS", size=7.5, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    txb(slide, l+0.18, t+2.00, w-0.36, 0.65,
        apres, size=9.5, color=BLEU_TEXTE, wrap=True)
    # KPI — bande bas
    rect(slide, l+0.18, t+2.78, w-0.36, 0.24, fill=BLEU_FONCE)
    txb(slide, l+0.22, t+2.79, w-0.44, 0.22,
        f"KPI  ▸  {kpi}", size=8, bold=True, color=BLANC)


def step_card(slide, l, t, w, num_label, date_label, desc):
    """
    Carte plan 90j — layout épuré, une seule zone de texte par info.
    num_label : '01', date_label : 'J+0', desc : texte court
    """
    h = 1.90
    rect(slide, l, t, w, h, fill=CREME)
    # Barre colorée top
    rect(slide, l, t, w, 0.07, fill=BLEU_FONCE)
    # Numéro — une seule ligne, pas de wrap, taille réduite
    txb(slide, l+0.14, t+0.14, 0.60, 0.38,
        num_label, size=22, bold=True, color=BLEU_FONCE,
        align=PP_ALIGN.LEFT, wrap=False)
    # Date — placée à droite du numéro, même ligne
    txb(slide, l+0.82, t+0.22, w-0.98, 0.28,
        date_label, size=10, bold=True, color=ORANGE,
        align=PP_ALIGN.LEFT, wrap=False)
    # Séparateur fin
    rect(slide, l+0.14, t+0.60, w-0.28, 0.02, fill=BLEU_GRIS)
    # Description — zone ample, pas de débordement
    txb(slide, l+0.14, t+0.70, w-0.28, 1.05,
        desc, size=9.5, color=BLEU_TEXTE, wrap=True)


def quick_win_row(slide, l, t, w, text, idx):
    y = t + idx * 0.62
    rect(slide, l, y, 0.26, 0.26, fill=TEAL)
    txb(slide, l, y, 0.26, 0.26,
        "✓", size=11, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
    txb(slide, l+0.36, y+0.02, w-0.36, 0.40,
        text, size=10.5, color=BLEU_TEXTE, wrap=True)


def principe_row(slide, l, t, w, text, idx):
    y = t + idx * 0.62
    txb(slide, l, y, w, 0.46,
        f"▸  {text}", size=10.5, color=BLANC, wrap=True)


# ══════════════════════════════════════════════════════════════════════════
# DONNÉES
# ══════════════════════════════════════════════════════════════════════════
PROSPECT     = "Orange Côte d'Ivoire"
SECTEUR      = "IT & Transformation Digitale"
TYPE_ORG     = "Entreprise de télécommunications"
PAYS         = "Côte d'Ivoire"
PACK_ID      = "pack-1780665567601-jp48u6z4"
AUDIT_URL    = f"https://www.transferai.ci/questionnaire-audit?pack_id={PACK_ID}"
CALENDLY_URL = "https://calendly.com/contact-transferai/30min"

CONTEXTE = (
    "Orange Côte d'Ivoire opère des flux à fort volume où la qualité de réponse, "
    "la coordination entre canaux et la diffusion des procédures impactent directement "
    "l'expérience client et l'efficacité interne."
)

ENJEUX = [
    "Amélioration des processus de reporting financier",
    "Formation du personnel sur les nouvelles technologies IA",
    "Optimisation des services clients multicanaux",
]

USE_CASES = [
    ("💬", "Copilote service client multicanal",
     "Demandes récurrentes réparties sur plusieurs canaux avec des réponses hétérogènes.",
     "Copilote de réponse unifié, résumés d'interactions et continuité améliorée.",
     "Délai moyen de réponse · Taux résolution 1er contact"),
    ("📋", "Synthèse tickets et incidents",
     "Analyse manuelle des tickets et reporting lent pour les managers.",
     "Synthèses automatiques, regroupement des motifs et alertes plus lisibles.",
     "Temps de préparation reporting · Délai d'escalade"),
    ("📖", "Base procédures et scripts terrain",
     "Réponses et procédures dispersées selon les équipes et les documents.",
     "Base de connaissances assistée pour scripts, contrôles et bonnes réponses.",
     "Temps recherche d'information · Temps onboarding"),
    ("📊", "Commentaire IA du pilotage qualité",
     "Les indicateurs existent mais restent peu commentés et difficiles à exploiter.",
     "Commentaires automatiques des écarts, synthèses managers et meilleurs arbitrages.",
     "Temps lecture des KPI · Réactivité de pilotage"),
]

QUICK_WINS = [
    "Réduction du délai moyen de réponse service client",
    "Accélération de la préparation des reportings managériaux",
    "Meilleure homogénéité des réponses terrain",
    "Réduction du temps d'onboarding des nouvelles recrues",
]

PRINCIPES_GOV = [
    "Supervision humaine de tous les usages IA",
    "Cadre de validation des sorties et des KPI",
    "Pilotage simple, sobre et explicable",
    "Arbitrage d'extension basé sur les résultats mesurés",
]

PLAN_90J = [
    ("01", "J+0",  "Audit & cadrage des flux service client et pilotage prioritaires."),
    ("02", "J+15", "Premier copilote ou première synthèse automatisée déployée."),
    ("03", "J+45", "Formation ciblée des équipes et lancement du pilote à l'échelle."),
    ("04", "J+90", "Mesure des résultats et décision d'extension sur de nouveaux flux."),
]

EQUIPE = [
    ("Casimir Kassi Beda",  "Directeur Général — Vision stratégique et pilotage",  "contact@transferai.ci"),
    ("Marius Ayoro",        "Directeur Développement et Partenariats",              "marius.ayoro@transferai.ci"),
    ("Eric N'Guessan",      "Directeur Audit, Pédagogie et Certification",          "eric.nguessan@transferai.ci"),
]


# ══════════════════════════════════════════════════════════════════════════
# CONSTRUCTION DES SLIDES
# ══════════════════════════════════════════════════════════════════════════
def build_deck():
    prs = new_prs()

    # ── SLIDE 1 — COUVERTURE ─────────────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 8.80, 0, 4.533, 7.5, fill=BLEU_MOYEN)
    rect(sl, 0.68, 0.90, 0.06, 5.20, fill=ORANGE)

    txb(sl, 1.00, 0.52, 3.80, 0.28,
        "DECK COMMERCIAL PREMIUM",
        size=8, bold=True, color=OR_AMBRE)
    txb(sl, 1.02, 1.00, 7.40, 0.65,
        f"Présentation TransferAI", size=26, bold=True, color=BLANC)
    txb(sl, 1.02, 1.72, 7.40, 0.48,
        f"pour {PROSPECT}", size=22, bold=True, color=ORANGE)
    txb(sl, 1.04, 2.42, 7.20, 1.20, CONTEXTE,
        size=11, color=BLEU_PALE, wrap=True)
    txb(sl, 1.04, 3.80, 7.0, 0.24,
        f"{SECTEUR}  ·  {TYPE_ORG}  ·  {PAYS}",
        size=8.5, color=BLEU_GRIS)
    pill_shape(sl, 1.02, 4.42, 2.30, 0.50, "Audit + exécution terrain",
               bg=BLEU_CLAIR, size=9)
    pill_shape(sl, 3.48, 4.42, 2.20, 0.50, "Formation & gouvernance",
               bg=ORANGE, size=9)

    # Colonne droite
    txb(sl, 9.10, 1.05, 3.50, 0.28,
        "Positionnement TransferAI × NettelecomCI",
        size=8, bold=True, color=OR_AMBRE)
    txb(sl, 9.10, 1.48, 3.45, 1.30,
        "Présence locale forte en Côte d'Ivoire, exécution internationale, "
        "méthodologie d'audit, formation ciblée et accompagnement sur 90 jours.",
        size=10, color=BLANC, wrap=True)
    txb(sl, 9.10, 3.10, 1.80, 0.24, "CTA prioritaire",
        size=8, bold=True, color=ORANGE)
    txb(sl, 9.10, 3.48, 3.45, 0.85,
        "Planifier un audit stratégique gratuit et un rendez-vous expert de 30 minutes",
        size=11, bold=True, color=BLANC, wrap=True)
    txb(sl, 9.10, 4.52, 1.10, 0.22, "Lien direct →",
        size=8, color=BLEU_PALE)
    txb(sl, 9.10, 4.80, 3.40, 0.32, CALENDLY_URL,
        size=8.5, color=BLEU_PALE, italic=True)
    footer(sl, 1)

    # ── SLIDE 2 — AGENDA ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)

    txb(sl, 0.78, 0.50, 6.0, 0.55,
        "Agenda de la présentation",
        size=22, bold=True, color=BLEU_FONCE)
    txb(sl, 0.80, 1.15, 6.5, 0.50,
        "Un fil narratif commercial premium pour relier enjeux métier, "
        "cas d'usage, ROI et accompagnement sur 90 jours.",
        size=10.5, color=BLEU_TEXTE, italic=True)

    agenda = [
        ("01", "Contexte & ambition",      "Positionnement, enjeux et objectif du deck"),
        ("02", "TransferAI & NettelecomCI", "Crédibilité, équipe et mode d'intervention"),
        ("03", "Enjeux du prospect",        "Pain points, besoins et signaux pré-audit"),
        ("04", "Cas d'usage prioritaires",  "Pistes de transformation IA ciblées"),
        ("05", "ROI & gouvernance",         "KPI, gains attendus et cadre de pilotage"),
        ("06", "Plan 90 jours",             "Formation, support et décision d'extension"),
    ]
    for i, (num, titre, sous) in enumerate(agenda):
        agenda_row(sl, num, titre, sous, y=1.90 + i * 0.74)

    # Encart droit
    rect(sl, 7.55, 1.80, 4.90, 3.90, fill=BLEU_FONCE)
    txb(sl, 7.85, 2.08, 4.30, 0.28, "Message central",
        size=9, bold=True, color=OR_AMBRE)
    txb(sl, 7.85, 2.50, 4.20, 1.15,
        "Transformer les priorités métier du prospect en gains visibles, "
        "mesurés et défendables dès les 90 premiers jours.",
        size=12, bold=True, color=BLANC, wrap=True)
    txb(sl, 7.85, 3.85, 4.20, 1.40,
        "Ce deck est un support commercial premium structuré à partir "
        "du pré-audit, du contexte sectoriel et des cas d'usage "
        "prioritaires du prospect.",
        size=9.5, color=BLEU_PALE, wrap=True)
    footer(sl, 2)

    # ── SLIDE 3 — ABOUT TRANSFERAI ────────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 0.68, 0.95, 0.06, 5.10, fill=ORANGE)

    txb(sl, 1.02, 1.00, 3.50, 0.28, "TransferAI × NettelecomCI",
        size=9, bold=True, color=OR_AMBRE)
    txb(sl, 1.02, 1.42, 6.20, 0.85,
        "Une exécution locale et internationale pour transformer "
        "l'IA en résultats concrets et mesurables.",
        size=18, bold=True, color=BLANC, wrap=True)
    txb(sl, 1.05, 2.50, 5.90, 0.95,
        "TransferAI est le hub IA de NettelecomCI en Côte d'Ivoire. "
        "Nous combinons présence locale, expertise sectorielle et "
        "exécution internationale pour déployer des cas d'usage utiles, "
        "gouvernés et mesurables.",
        size=10.5, color=BLEU_PALE, wrap=True)

    atouts = [
        ("🌍", "Présence locale forte via NettelecomCI"),
        ("👥", "13 experts ivoiriens et internationaux"),
        ("✈",  "Interventions CI, Royaume-Uni, USA, Inde"),
        ("🎯", "Formation + implémentation + gouvernance"),
    ]
    for i, (ico, txt) in enumerate(atouts):
        col, row = i % 2, i // 2
        lx = 1.05 + col * 3.10
        ty = 3.72 + row * 1.02
        rect(sl, lx, ty, 2.80, 0.84, fill=BLEU_MOYEN)
        txb(sl, lx+0.12, ty+0.14, 0.42, 0.56, ico, size=20, color=BLANC)
        txb(sl, lx+0.62, ty+0.18, 2.05, 0.55, txt,
            size=9.5, color=BLANC, wrap=True)

    # Colonne droite
    rect(sl, 7.72, 1.05, 4.85, 5.35, fill=BLEU_MOYEN)
    txb(sl, 8.02, 1.42, 4.30, 0.28, "Positionnement pour ce prospect",
        size=9, bold=True, color=OR_AMBRE)
    txb(sl, 8.02, 1.85, 4.10, 1.25, CONTEXTE,
        size=10.5, color=BLANC, wrap=True)
    txb(sl, 8.02, 3.32, 2.00, 0.26, "Promesse TransferAI",
        size=9, bold=True, color=ORANGE)
    txb(sl, 8.02, 3.68, 4.10, 0.88,
        "Audit ciblé, copilote service client, synthèse managériale "
        "et accompagnement sur 90 jours.",
        size=10.5, color=BLANC, wrap=True)
    txb(sl, 8.02, 4.75, 4.10, 0.80,
        CALENDLY_URL, size=8.5, color=BLEU_PALE, italic=True)
    footer(sl, 3)

    # ── SLIDE 4 — ENJEUX ─────────────────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)

    txb(sl, 0.78, 0.50, 9.0, 0.55,
        "Enjeux du prospect et priorités de transformation",
        size=22, bold=True, color=BLEU_FONCE)
    txb(sl, 0.80, 1.15, 6.60, 0.65,
        f"Le pré-audit montre des enjeux prioritaires en {SECTEUR} "
        "qui appellent une feuille de route IA ciblée, pilotable et visible.",
        size=10.5, color=BLEU_TEXTE, italic=True)

    for i, enjeu in enumerate(ENJEUX):
        y = 2.05 + i * 1.00
        rect(sl, 0.88, y, 0.06, 0.65, fill=ORANGE)
        rect(sl, 1.05, y, 6.05, 0.70, fill=CREME)
        txb(sl, 1.25, y+0.10, 5.70, 0.50,
            enjeu, size=13, bold=True, color=BLEU_FONCE)

    # Colonne droite
    rect(sl, 7.72, 1.50, 4.72, 4.65, fill=BLEU_FONCE)
    txb(sl, 8.02, 1.88, 4.12, 0.30,
        "Cadre de réponse TransferAI", size=9, bold=True, color=OR_AMBRE)
    txb(sl, 8.02, 2.30, 4.12, 0.80,
        "Des gains visibles, un pilote concret et un cadre d'extension "
        "sécurisé dès les 90 premiers jours.",
        size=10.5, color=BLANC, wrap=True)
    txb(sl, 8.02, 3.28, 1.80, 0.26, "Signal principal",
        size=8, bold=True, color=ORANGE)
    txb(sl, 8.02, 3.65, 4.12, 0.62,
        "Prioriser un point d'entrée concret, mesurable et défendable.",
        size=10.5, color=BLANC, wrap=True)
    txb(sl, 8.02, 4.45, 2.00, 0.26, "Cas d'usage recommandé",
        size=8, bold=True, color=ORANGE)
    txb(sl, 8.02, 4.82, 4.12, 0.62,
        "Structurer un premier cas d'usage IA à fort impact opérationnel.",
        size=10.5, color=BLANC, wrap=True)
    footer(sl, 4)

    # ── SLIDES 5 & 6 — CAS D'USAGE ───────────────────────────────────────
    for slide_idx in range(2):
        sl = blank_slide(prs)
        rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)
        txb(sl, 0.78, 0.48, 11.5, 0.55,
            "Cas d'usage prioritaires à présenter au prospect",
            size=22, bold=True, color=BLEU_FONCE)
        txb(sl, 0.80, 1.10, 11.5, 0.52,
            "Une sélection premium de cas d'usage transformant les enjeux "
            "du pré-audit en scénarios concrets de performance.",
            size=10.5, color=BLEU_TEXTE, italic=True)
        pair = USE_CASES[slide_idx*2 : slide_idx*2+2]
        positions = [0.82, 6.92]
        for lx, (ico, titre, avant, apres, kpi) in zip(positions, pair):
            use_case_card(sl, lx, 1.88, 5.70, ico, titre, avant, apres, kpi)
        footer(sl, 5 + slide_idx)

    # ── SLIDE 7 — ROI & GOUVERNANCE ───────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)
    txb(sl, 0.78, 0.50, 11.0, 0.55,
        "ROI, gouvernance et indicateurs défendables",
        size=22, bold=True, color=BLEU_FONCE)
    txb(sl, 0.80, 1.12, 11.2, 0.60,
        "Une lecture business des gains : productivité, qualité de service, "
        "accélération du reporting et meilleure exécution terrain.",
        size=10.5, color=BLEU_TEXTE, italic=True)

    # Quick wins
    rect(sl, 0.88, 2.05, 5.90, 3.72, fill=CREME)
    txb(sl, 1.10, 2.28, 5.40, 0.34,
        "Quick wins à objectiver dès la phase pilote",
        size=11, bold=True, color=BLEU_FONCE)
    for i, qw in enumerate(QUICK_WINS):
        quick_win_row(sl, 1.10, 2.82, 5.30, qw, i)

    # Gouvernance
    rect(sl, 7.08, 2.05, 5.35, 3.72, fill=BLEU_FONCE)
    txb(sl, 7.30, 2.28, 4.80, 0.34,
        "Principes de gouvernance", size=11, bold=True, color=OR_AMBRE)
    for i, p in enumerate(PRINCIPES_GOV):
        principe_row(sl, 7.30, 2.82, 4.80, p, i)
    footer(sl, 7)

    # ── SLIDE 8 — PLAN 90 JOURS ───────────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)
    txb(sl, 0.78, 0.50, 11.5, 0.55,
        "Trajectoire 90 jours : audit, formation, accompagnement",
        size=21, bold=True, color=BLEU_FONCE)
    txb(sl, 0.80, 1.12, 11.5, 0.55,
        "Un chemin premium standardisé mais adapté au prospect : cadrer, lancer, "
        "former, mesurer et décider l'extension sur des cas réels.",
        size=10.5, color=BLEU_TEXTE, italic=True)

    # 4 étapes — chacune dans step_card sans chevauchement
    for i, (num, date, desc) in enumerate(PLAN_90J):
        lx = 0.88 + i * 3.10
        step_card(sl, lx, 1.95, 2.82, num, date, desc)

    # Focus formation
    rect(sl, 0.88, 4.10, 5.75, 2.12, fill=CREME)
    txb(sl, 1.10, 4.28, 2.20, 0.32, "Focus formation",
        size=10, bold=True, color=BLEU_FONCE)
    formation = [
        "Usage quotidien du copilote de réponse",
        "Lecture et validation des synthèses IA",
        "Gouvernance des réponses et supervision humaine",
        "Mise à jour des scripts et procédures",
    ]
    for i, f in enumerate(formation):
        txb(sl, 1.10, 4.68 + i*0.34, 5.30, 0.32,
            f"•  {f}", size=9.5, color=BLEU_TEXTE)

    # Accompagnement
    rect(sl, 7.00, 4.10, 5.45, 2.12, fill=BLEU_FONCE)
    txb(sl, 7.22, 4.28, 4.80, 0.32, "Accompagnement 90 jours",
        size=10, bold=True, color=OR_AMBRE)
    accomp = [
        "Points de suivi d'usage et de pilotage",
        "Revue de cas réels et blocages terrain",
        "Ajustement des prompts, scripts et procédures",
        "Lecture des premiers KPI et arbitrage d'extension",
    ]
    for i, a in enumerate(accomp):
        txb(sl, 7.22, 4.68 + i*0.34, 4.90, 0.32,
            f"•  {a}", size=9.5, color=BLANC)
    footer(sl, 8)

    # ── SLIDE 9 — CTA FINAL ───────────────────────────────────────────────
    sl = blank_slide(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 7.18, 0, 6.153, 7.5, fill=BLEU_MOYEN)
    rect(sl, 0.68, 0.95, 0.06, 5.60, fill=ORANGE)

    txb(sl, 0.98, 1.05, 3.80, 0.30,
        "Prochaine étape recommandée", size=9, bold=True, color=OR_AMBRE)
    txb(sl, 0.98, 1.48, 6.00, 1.05,
        "Planifier un audit stratégique gratuit "
        "et un rendez-vous expert de 30 minutes",
        size=21, bold=True, color=BLANC, wrap=True)

    rect(sl, 0.98, 2.82, 5.90, 0.58, fill=BLEU_CLAIR)
    txb(sl, 1.18, 2.93, 1.30, 0.28, "Formulaire ▸",
        size=8, bold=True, color=OR_AMBRE)
    txb(sl, 2.60, 2.93, 4.00, 0.28, AUDIT_URL,
        size=8.5, color=BLEU_PALE, italic=True)

    rect(sl, 0.98, 3.58, 5.90, 0.58, fill=BLEU_CLAIR)
    txb(sl, 1.18, 3.69, 1.40, 0.28, "Rendez-vous ▸",
        size=8, bold=True, color=OR_AMBRE)
    txb(sl, 2.72, 3.69, 3.88, 0.28, CALENDLY_URL,
        size=8.5, color=BLEU_PALE, italic=True)

    txb(sl, 0.98, 4.42, 5.90, 1.15,
        "Un dispositif premium pour convaincre, exécuter vite "
        "et mesurer les gains dès les 90 premiers jours.",
        size=12, bold=True, color=BLANC, italic=True, wrap=True)

    # Équipe
    txb(sl, 7.48, 1.45, 4.50, 0.32,
        "Équipe dirigeante mobilisable", size=10, bold=True, color=OR_AMBRE)
    for i, (nom, role, email) in enumerate(EQUIPE):
        y = 2.00 + i * 1.18
        txb(sl, 7.48, y,       4.50, 0.30, nom,
            size=11, bold=True, color=BLANC)
        txb(sl, 7.48, y+0.34,  4.50, 0.30, role,
            size=9, color=BLEU_PALE, wrap=True)
        txb(sl, 7.48, y+0.68,  4.50, 0.26, email,
            size=8.5, color=ORANGE, italic=True)
    footer(sl, 9)

    prs.save(OUT_PPTX)
    print(f"✓ PPTX : {OUT_PPTX}")
    return OUT_PPTX


def export_pdf(pptx_path):
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    script = f'''
tell application "Keynote"
    set theDoc to open POSIX file "{pptx_path}"
    delay 4
    export theDoc to POSIX file "{pdf_path}" as PDF with properties {{PDF image quality:Best}}
    close theDoc saving no
end tell
'''
    try:
        r = subprocess.run(["osascript", "-e", script],
                           capture_output=True, text=True, timeout=90)
        if r.returncode == 0:
            print(f"✓ PDF  : {pdf_path}")
        else:
            print(f"⚠ PDF auto échoué : {r.stderr.strip()}")
    except Exception as e:
        print(f"⚠ PDF : {e}")
    return pdf_path


if __name__ == "__main__":
    pptx = build_deck()
    export_pdf(pptx)
    print("✓ Terminé — fichiers dans ~/Downloads/")
