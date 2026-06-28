#!/usr/bin/env python3
"""
deck_generator_cli.py — Générateur de deck TransferAI
Appelé par n8n via Execute Command.

Usage:
  python3 deck_generator_cli.py '<JSON>'

JSON attendu:
{
  "prospect": "Orange Côte d'Ivoire",
  "secteur": "Télécommunications",
  "pays": "Côte d'Ivoire",
  "pack_id": "pack-xxx",
  "accroche": "Transformer vos flux clients avec l'IA",
  "enjeux": ["enjeu 1", "enjeu 2", "enjeu 3"],
  "cas_usage": [
    ["💬", "Titre cas 1", "Gain attendu 1"],
    ["📋", "Titre cas 2", "Gain attendu 2"],
    ["📖", "Titre cas 3", "Gain attendu 3"]
  ],
  "gains": [
    ["Délai de réponse", "− 30 à 50 %", "teal"],
    ["Temps de reporting", "− 60 %", "orange"],
    ["Qualité réponses", "+ homogénéité", "blue"]
  ],
  "supabase_url": "https://wlhznciwuofueffyoflo.supabase.co",
  "supabase_service_key": "eyJ..."
}

Sortie (stdout JSON):
{
  "ok": true,
  "pptx_url": "https://...supabase.co/storage/v1/object/public/prospect-decks/deck_xxx.pptx",
  "pdf_url":  "https://...supabase.co/storage/v1/object/public/prospect-decks/deck_xxx.pdf",
  "filename_pptx": "Deck_TransferAI_OrangeCoteIvoire.pptx",
  "filename_pdf":  "Deck_TransferAI_OrangeCoteIvoire.pdf"
}
"""

import sys
import json
import os
import re
import tempfile
import subprocess
import urllib.request

# ── Dépendance : python-pptx ──────────────────────────────────────────
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print(json.dumps({"ok": False, "error": "python-pptx non installé. Lancer: pip3 install python-pptx requests"}))
    sys.exit(1)


# ── Palette ───────────────────────────────────────────────────────────
BLEU_FONCE = RGBColor(0x10, 0x26, 0x3F)
BLEU_MOYEN = RGBColor(0x16, 0x35, 0x56)
BLEU_CLAIR = RGBColor(0x24, 0x46, 0x67)
BLEU_TEXTE = RGBColor(0x61, 0x70, 0x86)
ORANGE     = RGBColor(0xE7, 0x6F, 0x1D)
OR_AMBRE   = RGBColor(0xC8, 0x8C, 0x3A)
TEAL       = RGBColor(0x1C, 0x8A, 0x78)
BLANC      = RGBColor(0xFF, 0xFF, 0xFF)
CREME      = RGBColor(0xF5, 0xEF, 0xE7)
BLEU_PALE  = RGBColor(0xCF, 0xE0, 0xF0)
BLEU_GRIS  = RGBColor(0xC9, 0xD4, 0xE1)

COLOR_MAP = {
    "teal":   TEAL,
    "orange": ORANGE,
    "blue":   BLEU_CLAIR,
    "or":     OR_AMBRE,
}

EQUIPE = [
    ("Casimir Kassi Beda", "Directeur Général",                          "contact@transferai.ci"),
    ("Marius Ayoro",       "Directeur Développement & Partenariats",     "marius.ayoro@transferai.ci"),
    ("Soulemane Konate",   "Directeur IA & Innovation",                  "soulemane.konate@transferai.ci"),
    ("Eric N'Guessan",     "Directeur Audit, Pédagogie & Certification", "eric.nguessan@transferai.ci"),
    ("Médard Séry",        "Consultant expert · Data & plateformes IA",  "medard.sery@transferai.ci"),
]

PLAN_90J_DEFAULT = [
    ("01", "J+0",  "Audit & cadrage"),
    ("02", "J+15", "Premier pilote"),
    ("03", "J+45", "Formation & déploiement"),
    ("04", "J+90", "Mesure & extension"),
]


# ── Helpers PPTX ─────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s

def tx(slide, l, t, w, h, text, size=11, bold=False, color=BLANC,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tf = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf.text_frame.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = str(text)
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tf

def footer(slide, num):
    rect(slide, 0.5, 7.02, 12.33, 0.02, fill=ORANGE)
    tx(slide, 0.6, 7.06, 10.0, 0.22,
       "TransferAI  ·  Hub IA de NettelecomCI  ·  contact@transferai.ci",
       size=7, color=BLEU_TEXTE)
    tx(slide, 12.2, 7.04, 0.80, 0.24,
       f"{num:02d}", size=9, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)


# ── Construction slides ───────────────────────────────────────────────
def build_deck(d):
    prospect  = d.get("prospect", "Votre organisation")
    secteur   = d.get("secteur", "")
    pays      = d.get("pays", "")
    pack_id   = d.get("pack_id", "")
    accroche  = d.get("accroche", "Transformer vos processus avec l'IA")
    enjeux    = d.get("enjeux", ["Enjeu 1", "Enjeu 2", "Enjeu 3"])[:3]
    cas_usage = d.get("cas_usage", [["💡","Cas d'usage 1","Gain attendu"]])[:3]
    gains_raw = d.get("gains", [["KPI 1","− 30 %","teal"]])[:3]
    plan      = d.get("plan_90j", PLAN_90J_DEFAULT)[:4]

    audit_url    = f"https://www.transferai.ci/questionnaire-audit?pack_id={pack_id}" if pack_id else "https://www.transferai.ci/questionnaire-audit"
    calendly_url = "https://calendly.com/contact-transferai/30min"

    gains = []
    for g in gains_raw:
        label, valeur = g[0], g[1]
        col_key = g[2] if len(g) > 2 else "teal"
        gains.append((label, valeur, COLOR_MAP.get(col_key, TEAL)))

    prs = new_prs()

    # ── SLIDE 1 — Couverture ─────────────────────────────────────────
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 8.60, 0, 4.733, 7.5, fill=BLEU_MOYEN)
    rect(sl, 0.55, 0.85, 0.05, 5.50, fill=ORANGE)
    tx(sl, 0.88, 0.82, 4.0, 0.26, "PRÉSENTATION COMMERCIALE",
       size=7.5, bold=True, color=OR_AMBRE)
    tx(sl, 0.90, 1.28, 7.40, 0.72, f"TransferAI × {prospect}",
       size=24, bold=True, color=BLANC)
    tx(sl, 0.90, 2.16, 7.30, 0.38, accroche,
       size=13, color=BLEU_PALE, italic=True)
    rect(sl, 0.90, 2.76, 3.20, 0.03, fill=ORANGE)
    tx(sl, 0.90, 2.98, 7.0, 0.28, f"{secteur}  ·  {pays}", size=10, color=BLEU_GRIS)
    for i, e in enumerate(enjeux):
        tx(sl, 0.90, 3.55 + i * 0.52, 7.20, 0.42, f"▸  {e}", size=10.5, color=BLANC)
    tx(sl, 8.90, 1.20, 3.80, 0.28, "TransferAI Africa", size=9, bold=True, color=OR_AMBRE)
    tx(sl, 8.90, 1.62, 3.70, 0.26, "Hub IA de NettelecomCI", size=10.5, bold=True, color=BLANC)
    tx(sl, 8.90, 2.06, 3.70, 1.00,
       "Présence locale CI · Expertise internationale · Audit, formation et accompagnement 90 jours.",
       size=9.5, color=BLEU_PALE, wrap=True)
    rect(sl, 8.90, 3.22, 3.70, 0.02, fill=BLEU_GRIS)
    tx(sl, 8.90, 3.40, 3.70, 0.26, "Rendez-vous →", size=9, bold=True, color=ORANGE)
    tx(sl, 8.90, 3.76, 3.70, 0.28, calendly_url, size=8, color=BLEU_PALE, italic=True)
    tx(sl, 8.90, 4.24, 3.70, 0.26, "Formulaire audit →", size=9, bold=True, color=ORANGE)
    tx(sl, 8.90, 4.60, 3.70, 0.28, "transferai.ci/questionnaire-audit",
       size=8, color=BLEU_PALE, italic=True)
    footer(sl, 1)

    # ── SLIDE 2 — Enjeux ────────────────────────────────────────────
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)
    rect(sl, 0, 0, 13.333, 1.40, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.30, 10.0, 0.50, "Enjeux identifiés lors du pré-audit",
       size=20, bold=True, color=BLANC)
    tx(sl, 0.74, 0.88, 9.0, 0.30, f"{prospect}  ·  {secteur}", size=9.5, color=BLEU_PALE)
    for i, enjeu in enumerate(enjeux):
        lx = 0.72 + i * 4.22
        rect(sl, lx, 1.75, 3.90, 3.20, fill=CREME)
        rect(sl, lx, 1.75, 3.90, 0.07, fill=ORANGE)
        tx(sl, lx+0.20, 1.90, 0.60, 0.55, f"0{i+1}", size=26, bold=True, color=ORANGE)
        tx(sl, lx+0.20, 2.62, 3.50, 1.80, enjeu, size=13, bold=True, color=BLEU_FONCE, wrap=True)
    rect(sl, 0.72, 5.20, 11.89, 0.82, fill=BLEU_FONCE)
    tx(sl, 1.00, 5.42, 11.40, 0.42,
       "Ces enjeux sont les points d'entrée prioritaires pour un premier déploiement IA ciblé, "
       "mesurable et défendable dès les 90 premiers jours.",
       size=10.5, color=BLANC, italic=True, wrap=True)
    footer(sl, 2)

    # ── SLIDE 3 — Cas d'usage ────────────────────────────────────────
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)
    rect(sl, 0, 0, 13.333, 1.40, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.30, 10.0, 0.50, "Cas d'usage prioritaires",
       size=20, bold=True, color=BLANC)
    tx(sl, 0.74, 0.88, 9.0, 0.30,
       "Trois leviers de transformation IA adaptés à votre secteur",
       size=9.5, color=BLEU_PALE)
    for i, cu in enumerate(cas_usage):
        ico, titre, gain = cu[0], cu[1], cu[2]
        lx = 0.72 + i * 4.22
        rect(sl, lx, 1.72, 3.90, 4.60, fill=CREME)
        rect(sl, lx, 1.72, 3.90, 0.07, fill=BLEU_FONCE)
        rect(sl, lx+0.20, 1.95, 0.70, 0.70, fill=BLEU_FONCE)
        tx(sl, lx+0.20, 1.92, 0.70, 0.70, ico, size=22, color=BLANC, align=PP_ALIGN.CENTER)
        tx(sl, lx+1.05, 2.05, 2.65, 0.55, titre, size=13, bold=True, color=BLEU_FONCE, wrap=True)
        rect(sl, lx+0.20, 2.82, 3.50, 0.02, fill=BLEU_GRIS)
        rect(sl, lx+0.20, 3.06, 0.65, 0.24, fill=TEAL)
        tx(sl, lx+0.20, 3.06, 0.65, 0.24, "GAIN", size=7, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
        tx(sl, lx+0.20, 3.42, 3.50, 1.60, gain, size=11, color=BLEU_TEXTE, wrap=True)
    footer(sl, 3)

    # ── SLIDE 4 — ROI ────────────────────────────────────────────────
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.65, 10.0, 0.55, "Résultats attendus dès le pilote",
       size=22, bold=True, color=BLANC)
    tx(sl, 0.74, 1.30, 9.0, 0.30, "Des gains visibles, mesurés et défendables",
       size=10.5, color=BLEU_PALE, italic=True)
    for i, (label, valeur, couleur) in enumerate(gains):
        lx = 0.72 + i * 4.22
        rect(sl, lx, 2.00, 3.90, 2.20, fill=BLEU_MOYEN)
        rect(sl, lx, 2.00, 3.90, 0.07, fill=couleur)
        tx(sl, lx+0.24, 2.28, 3.42, 0.40, label, size=10, color=BLEU_PALE)
        tx(sl, lx+0.24, 2.80, 3.42, 0.70, valeur, size=28, bold=True, color=couleur, wrap=False)
    rect(sl, 0.72, 4.52, 11.89, 1.60, fill=BLEU_MOYEN)
    tx(sl, 1.00, 4.75, 3.0, 0.30, "Gouvernance & supervision", size=9, bold=True, color=OR_AMBRE)
    principes = [
        "Supervision humaine de tous les usages IA",
        "KPI validés avant toute extension",
        "Pilotage simple, sobre et explicable",
    ]
    for i, p in enumerate(principes):
        tx(sl, 1.00 + i * 3.96, 5.15, 3.70, 0.70, f"▸  {p}", size=10, color=BLANC, wrap=True)
    footer(sl, 4)

    # ── SLIDE 5 — Plan 90 jours ──────────────────────────────────────
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLANC)
    rect(sl, 0, 0, 13.333, 1.40, fill=BLEU_FONCE)
    tx(sl, 0.72, 0.30, 10.0, 0.50, "Plan d'action 90 jours",
       size=20, bold=True, color=BLANC)
    tx(sl, 0.74, 0.88, 9.0, 0.30, "Audit · Formation · Déploiement · Extension",
       size=9.5, color=BLEU_PALE)
    for i, (num, delai, action) in enumerate(plan):
        lx = 0.72 + i * 3.15
        rect(sl, lx, 1.72, 2.90, 2.30, fill=CREME)
        rect(sl, lx, 1.72, 2.90, 0.07, fill=BLEU_FONCE)
        tx(sl, lx+0.20, 1.88, 0.70, 0.50, num, size=26, bold=True, color=BLEU_FONCE, wrap=False)
        tx(sl, lx+1.00, 1.96, 1.70, 0.34, delai, size=13, bold=True, color=ORANGE, wrap=False)
        rect(sl, lx+0.20, 2.52, 2.50, 0.02, fill=BLEU_GRIS)
        tx(sl, lx+0.20, 2.66, 2.50, 0.90, action, size=12, bold=True, color=BLEU_FONCE, wrap=True)
    rect(sl, 0.72, 4.26, 5.75, 2.00, fill=BLEU_FONCE)
    tx(sl, 0.96, 4.48, 5.20, 0.30, "Accompagnement TransferAI", size=9, bold=True, color=OR_AMBRE)
    for i, a in enumerate(["Suivi hebdomadaire sur l'usage réel",
                            "Ajustement des prompts et procédures",
                            "Revue des KPI et décision d'extension"]):
        tx(sl, 0.96, 4.90 + i * 0.38, 5.20, 0.34, f"·  {a}", size=10, color=BLANC)
    rect(sl, 6.72, 4.26, 5.89, 2.00, fill=CREME)
    tx(sl, 6.96, 4.48, 5.40, 0.30, "Formation ciblée", size=9, bold=True, color=BLEU_FONCE)
    for i, f in enumerate(["Usage du copilote en situation réelle",
                            "Validation et gouvernance des réponses IA",
                            "Mise à jour des scripts et procédures"]):
        tx(sl, 6.96, 4.90 + i * 0.38, 5.40, 0.34, f"·  {f}", size=10, color=BLEU_TEXTE)
    footer(sl, 5)

    # ── SLIDE 6 — CTA + Équipe ───────────────────────────────────────
    sl = blank(prs)
    rect(sl, 0, 0, 13.333, 7.5, fill=BLEU_FONCE)
    rect(sl, 7.80, 0, 5.533, 7.5, fill=BLEU_MOYEN)
    rect(sl, 0.55, 0.85, 0.05, 5.50, fill=ORANGE)
    tx(sl, 0.88, 0.90, 4.0, 0.28, "PROCHAINE ÉTAPE", size=7.5, bold=True, color=OR_AMBRE)
    tx(sl, 0.88, 1.35, 6.60, 0.85,
       "Planifier un audit stratégique gratuit\net un rendez-vous expert de 30 minutes",
       size=19, bold=True, color=BLANC, wrap=True)
    rect(sl, 0.88, 2.50, 6.55, 0.02, fill=ORANGE)
    rect(sl, 0.88, 2.72, 6.55, 0.64, fill=BLEU_CLAIR)
    tx(sl, 1.08, 2.84, 1.50, 0.26, "Rendez-vous", size=8.5, bold=True, color=OR_AMBRE)
    tx(sl, 2.72, 2.88, 4.50, 0.26, calendly_url, size=8.5, color=BLEU_PALE, italic=True)
    rect(sl, 0.88, 3.50, 6.55, 0.64, fill=BLEU_CLAIR)
    tx(sl, 1.08, 3.62, 1.70, 0.26, "Formulaire audit", size=8.5, bold=True, color=OR_AMBRE)
    tx(sl, 2.72, 3.66, 4.50, 0.26, "transferai.ci/questionnaire-audit",
       size=8.5, color=BLEU_PALE, italic=True)
    tx(sl, 0.88, 4.44, 6.55, 0.80,
       "Un premier échange de 30 minutes pour cadrer les priorités et définir le périmètre du pilote.",
       size=10.5, color=BLANC, italic=True, wrap=True)
    tx(sl, 8.10, 1.00, 4.80, 0.30, "Équipe dirigeante", size=9, bold=True, color=OR_AMBRE)
    for i, (nom, role, email) in enumerate(EQUIPE):
        y = 1.50 + i * 1.06
        rect(sl, 8.10, y, 4.80, 0.02, fill=BLEU_CLAIR)
        tx(sl, 8.10, y+0.10, 4.80, 0.28, nom, size=10.5, bold=True, color=BLANC)
        tx(sl, 8.10, y+0.40, 4.80, 0.26, role, size=8.5, color=BLEU_PALE, wrap=True)
        tx(sl, 8.10, y+0.68, 4.80, 0.22, email, size=7.5, color=ORANGE, italic=True)
    footer(sl, 6)

    return prs


# ── Conversion PDF (LibreOffice sur VPS Linux) ────────────────────────
def convert_to_pdf(pptx_path):
    out_dir = os.path.dirname(pptx_path)
    try:
        for cmd in ["libreoffice", "libreoffice7.6", "soffice"]:
            result = subprocess.run(
                [cmd, "--headless", "--convert-to", "pdf",
                 "--outdir", out_dir, pptx_path],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0:
                pdf_path = pptx_path.replace(".pptx", ".pdf")
                if os.path.exists(pdf_path):
                    return pdf_path
    except Exception:
        pass
    return None


# ── Upload Supabase Storage ───────────────────────────────────────────
def upload_to_supabase(file_path, filename, supabase_url, service_key):
    bucket = "prospect-decks"
    url = f"{supabase_url}/storage/v1/object/{bucket}/{filename}"
    content_type = "application/pdf" if filename.endswith(".pdf") else \
                   "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    with open(file_path, "rb") as f:
        data = f.read()
    req = urllib.request.Request(
        url,
        data=data,
        method="POST",
        headers={
            "Authorization": f"Bearer {service_key}",
            "Content-Type": content_type,
            "x-upsert": "true",
        }
    )
    try:
        with urllib.request.urlopen(req) as resp:
            resp.read()
    except urllib.error.HTTPError as e:
        body = e.read().decode()
        raise RuntimeError(f"Upload échoué ({e.code}): {body}")
    public_url = f"{supabase_url}/storage/v1/object/public/{bucket}/{filename}"
    return public_url


# ── Normalise nom de fichier ──────────────────────────────────────────
def safe_name(value):
    import unicodedata
    value = unicodedata.normalize("NFD", str(value))
    value = "".join(c for c in value if unicodedata.category(c) != "Mn")
    value = re.sub(r"[^a-zA-Z0-9]+", "_", value).strip("_")
    return value


# ── Main ──────────────────────────────────────────────────────────────
def main():
    if len(sys.argv) < 2:
        print(json.dumps({"ok": False, "error": "Argument JSON manquant"}))
        sys.exit(1)

    try:
        d = json.loads(sys.argv[1])
    except json.JSONDecodeError as e:
        print(json.dumps({"ok": False, "error": f"JSON invalide: {e}"}))
        sys.exit(1)

    supabase_url = d.get("supabase_url", "")
    service_key  = d.get("supabase_service_key", "")

    # Mode test local : génère le PPTX sans upload
    dry_run = (not supabase_url or supabase_url == "SKIP" or not service_key or service_key == "SKIP")
    if not supabase_url or not service_key:
        dry_run = True

    org_safe  = safe_name(d.get("prospect", "Prospect"))
    pack_id   = d.get("pack_id", "pack-unknown")
    stem      = f"Deck_TransferAI_{org_safe}_{pack_id}"
    fn_pptx   = stem + ".pptx"
    fn_pdf    = stem + ".pdf"

    with tempfile.TemporaryDirectory() as tmp:
        pptx_path = os.path.join(tmp, fn_pptx)
        prs = build_deck(d)
        prs.save(pptx_path)

        if dry_run:
            print(json.dumps({
                "ok": True,
                "dry_run": True,
                "pptx_generated": True,
                "pptx_url": None,
                "pdf_url": None,
                "filename_pptx": fn_pptx,
                "filename_pdf":  fn_pdf,
                "deck_artifact": None,
                "note": "Mode test local — upload Supabase ignoré (supabase_url ou service_key manquant)"
            }))
            return

        # Upload PPTX
        pptx_url = upload_to_supabase(pptx_path, fn_pptx, supabase_url, service_key)

        # Conversion PDF
        pdf_url = None
        pdf_path = convert_to_pdf(pptx_path)
        if pdf_path and os.path.exists(pdf_path):
            pdf_url = upload_to_supabase(pdf_path, fn_pdf, supabase_url, service_key)

    result = {
        "ok":            True,
        "pptx_url":      pptx_url,
        "pdf_url":       pdf_url,
        "filename_pptx": fn_pptx,
        "filename_pdf":  fn_pdf,
        "deck_artifact": {
            "pptx_url":      pptx_url,
            "pdf_url":       pdf_url,
            "filename_pptx": fn_pptx,
            "filename_pdf":  fn_pdf,
        }
    }
    print(json.dumps(result))


if __name__ == "__main__":
    main()
